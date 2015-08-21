#!/usr/bin/env python
import argparse
import tempfile
import zipfile
import glob
import pdb ### Remove
import platform
import sys
import os
import re
from PIL import ImageFont

from itertools import groupby
from operator import itemgetter
from xml.dom.minidom import parse
from shutil import rmtree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE


FONTSIZE=18

def getnotes(pptxfile):
    words = {}
    tmpd = tempfile.mkdtemp()
    zipfile.ZipFile(pptxfile).extractall(path=tmpd, pwd=None)

    # Parse notes content
    path = tmpd + '/ppt/notesSlides/'
    for infile in glob.glob(os.path.join(path, '*.xml')):
        #parse each XML notes file from the notes folder.
        dom = parse(infile)
        noteslist = dom.getElementsByTagName('a:t')

        # The page number is part of the filename
        page = int(re.sub(r'\D', "", infile.split("/")[-1]))

        text = ''
        for node in noteslist:
            xmlTag = node.toxml()
            xmlData = xmlTag.replace('<a:t>', '').replace('</a:t>', '')
            text += xmlData

        # Convert to ascii to simplify
        text = text.encode('ascii', 'ignore')
        words[page] = text

    # Remove all the files created with unzip
    rmtree(tmpd)
    return words

def createtoc(args, toc):

    # Create a new Presentation object, using the Style document as the template
    # The style document should be empty; otherwise we'll add the new ToC to the
    # end of the specified document.
    try:
        prs = Presentation(args.stylepptx)
    except:
        sys.stderr.write("Cannot read input style PowerPoint file \'%s\'. Possible malformed file.\n"%args.stylepptx)
        return

    # Create a blank slide in the object using the second master slide style by default
    blank_slide_layout = prs.slide_layouts[args.stylemasterslide]
    slide = prs.slides.add_slide(blank_slide_layout)

    # XXX TODO: Add handling when the number of ToC entries surpasses one slide
    slide.shapes.title.text = "Table of Contents"

    # Get font information
    font = ImageFont.truetype(args.fontdir + args.font, FONTSIZE)

    # This is the number of dots that fit using the given font in a 8.5" text box
    MAXDOTS=109
    # This is the maximum pixel width of a 8.5" text box
    MAXPXWIDTHTEXT=570
    MAXPXWIDTHBULLETS=580

    # This is the size of a single dot in pixels
    dotwidth=float(MAXPXWIDTHTEXT)/MAXDOTS

    # This is the maximum width for a given ToC line

    # The ToC entries and the page numbers are strings delimited by \n
    titles=''
    pages=''
    for pagenum in sorted(toc):
        tocpxlen = font.getsize(toc[pagenum])[0]
        if tocpxlen > MAXPXWIDTHTEXT:
            sys.stderr.write("Text for ToC entry on page %d (\"%s\") is too long, truncating.\n"%(pagenum, toc[pagenum]))
            # Trim one character off at a time until it fits! Presumably, the author will want to go back
            # and fix their original content for a smarter summarization of the ToC entry.
            while tocpxlen > MAXPXWIDTHTEXT:
                toc[pagenum] = toc[pagenum][:-1]
                tocpxlen = font.getsize(toc[pagenum])[0]
        titles += toc[pagenum] + "\n"
        pages += str(pagenum) + "\n"

    # Build the left-hand ToC entries first
    top=Inches(1.75)
    left=Inches(.5)
    width=Inches(8.5)
    height=Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size=None
    tf.text = titles
    p=tf.paragraphs[0]
    p.font.name = 'Tahoma'
    p.font.size=Pt(FONTSIZE)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    # Iterate through each of the ToC entries, calculating the number of dots needed in the middle textbox
    for page in sorted(toc):
        tocpxlen = font.getsize(toc[page])[0]
        #print "DEBUG: %03d %s"%(tocpxlen, toc[page])

        # The number of dots we use is the max width in pixels, minus the length of the ToC entry in pixels,
        # divided by the pixel width of a single dot, rounded down.
        tf.text+=("." * int(( (float(MAXPXWIDTHBULLETS - tocpxlen)) / dotwidth ))) + "\n"

    tf.auto_size=None
    p=tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = 'Tahoma'
    p.font.size=Pt(FONTSIZE)
    
    left=Inches(9)
    width=Inches(.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size=None
    tf.text = pages
    p=tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = 'Tahoma'
    p.font.size=Pt(FONTSIZE)
    
    try:
        prs.save(args.outputpptx)
    except:
        sys.stderr.write("Error saving output pptx file \'%s\'.\n"%args.outputpptx)
    return

class HelpWithRawFormatter(argparse.RawDescriptionHelpFormatter, argparse.ArgumentDefaultsHelpFormatter):
        pass

if __name__ == "__main__":

    # Establish a default for the location of fonts to make it easier to specify the location
    # of the font.
    deffontdir = '' # Works OK for Windows
    if platform.system() == 'Darwin':
        deffontdir='/Library/Fonts/'
    elif platform.system() == 'Linux':
        # Linux stores fonts in sub-directories, so users must specify sub-dir and font name
        # beneath this directory.
        deffontdir='/usr/share/fonts/truetype/'

    # Parse command-line arguments
    parser = argparse.ArgumentParser(
            description='Create a Table of Content from a PowerPoint file',
            prog='pptxtoc.py',
            formatter_class=HelpWithRawFormatter)
    parser.add_argument('-o', action="store", dest="outputpptx", default='toc.pptx',
            help="output pptx ToC slide")
    parser.add_argument('-f', action="store", dest="font", default='Tahoma',
            help="font filename")
    parser.add_argument('-F', action="store", dest="fontdir", default=deffontdir,
            help="font directory")
    parser.add_argument('-s', action="store", dest="stylepptx", default="Style.pptx",
            help="PowerPoint style document with no slides")
    parser.add_argument('-S', action="store", dest="stylemasterslide", type=int, default=2,
            help="slide number in master view to use for output ToC slide")
#    parser.add_argument('-z', action="store", dest="fontsize", type=int, default=FONTSIZE,
#            help="font size in points/pt")
    parser.add_argument('pptxfile')
    args = parser.parse_args()

    # Add the filename extension to font, relieving the user of this burden
    args.font += '.ttf'

    # Make sure the files exist and are readable
    if not (os.path.isfile(args.pptxfile) and os.access(args.pptxfile, os.R_OK)):
        sys.stderr.write("Cannot read the PowerPoint file \'%s\'.\n"%args.pptxfile)
        sys.exit(1)
    if not (os.path.isfile(args.stylepptx) and os.access(args.stylepptx, os.R_OK)):
        sys.stderr.write("Cannot read the PowerPoint style file \'%s\'.\n"%args.stylepptx)
        sys.exit(1)
    if not (os.path.isfile(args.fontdir + args.font) and os.access(args.fontdir + args.font, os.R_OK)):
        sys.stderr.write("Cannot read the the font file \'%s\'.\n"%(args.fontdir + args.font))
        sys.exit(1)

    # Decrement the style master slide number for offset counting
    if (args.stylemasterslide < 1):
        sys.stderr.write("Invalid style document master slide number \'%d\'.\n"%args.stylemasterslide)
        sys.exit(1)
    args.stylemasterslide -= 1


    # Retrieve all the notes from the pptx file in a page number-keyed dictionary
    words = getnotes(args.pptxfile)

    # Search for {{{whatever}}} and build a new dictionary of the page numbers and whatevers
    toc = {}
    for key in words:
        m=re.search(r'{{{(.*)}}}',words[key])
        if m is not None:
            toc[key] = m.groups(0)[0]

    # Generate the output ToC slide using the identified page numbers and titles
    createtoc(args, toc)

    print "Finished."
    sys.exit(0)
    
